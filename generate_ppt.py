"""
课程评价系统 - 演示PPT生成
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0, 51, 102)
BODY_COLOR = RGBColor(51, 51, 51)
ACCENT_COLOR = RGBColor(0, 102, 153)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # 顶部装饰条
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题栏
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.15)
    
    # 内容
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(10)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
    
    return slide

# ===== Slide 1: 封面 =====
add_title_slide(prs, '课程评价系统', '人工智能通识课·实践项目\n第一类：智能体编程应用开发')

# ===== Slide 2: 项目背景 =====
add_content_slide(prs, '项目背景与目标', [
    '● 课程评价是高校教学质量保障的重要环节，传统纸质评教效率低、统计难',
    '● 本项目开发一套线上课程评价系统，实现评价的实时采集与自动统计',
    '● 核心目标：构建轻量级、易部署的Web应用，满足日常教学评价需求',
    '● 技术栈：Python Flask + SQLite + HTML/CSS/JS + ECharts',
])

# ===== Slide 3: 系统角色 =====
add_content_slide(prs, '角色设计与权限分流', [
    '● 学生（Student）：浏览课程、选课、提交评价（1-5星+文字评语）、查看历史',
    '● 教师（Teacher）：查看评分统计与分布、浏览学生详细评价、导出CSV',
    '● 管理员（Admin）：课程增删改、用户角色管理、全校数据仪表板',
    '● 权限控制：基于RBAC模型，不同角色登录后显示完全不同的功能界面',
    '● 后端API层使用装饰器进行角色校验，防止越权访问',
])

# ===== Slide 4: 技术架构 =====
add_content_slide(prs, '技术架构与数据库设计', [
    '● B/S三层架构：表现层（HTML/JS）→ 业务逻辑层（Flask）→ 数据层（SQLite）',
    '● 4张核心数据表：users / courses / enrollments / evaluations',
    '● 前端通过Fetch API异步调用后端RESTful接口，页面无刷新切换内容',
    '● Session-based认证：登录后生成session，各页面加载时验证身份',
    '● ECharts CDN集成，实时查询数据库生成可交互图表',
])

# ===== Slide 5: 功能演示 =====
add_content_slide(prs, '功能演示（截图）', [
    '（此处插入登录页截图）',
    '（此处插入学生端截图 — 选课、评价、历史）',
    '（此处插入教师端截图 — 评分统计、评价列表、CSV导出）',
    '（此处插入管理员端截图 — 课程管理、用户管理、数据仪表板）',
])

# ===== Slide 6: 数据可视化 =====
add_content_slide(prs, '数据可视化（ECharts）', [
    '● 课程平均评分排名柱状图 — 直观比较各课程评价水平',
    '● 评分分布饼图 — 展示1-5星评价的数量和占比',
    '● 选课人数对比图 — 各课程热度一目了然',
    '● 所有图表支持交互式tooltip悬停查看详情',
    '（此处插入ECharts图表截图）',
])

# ===== Slide 7: 代码管理与部署 =====
add_content_slide(prs, 'GitHub仓库与代码管理', [
    '● 仓库地址：https://github.com/cty23/project（公开）',
    '● 完整项目结构：backend / frontend / database / README / requirements',
    '● 包含完整运行说明：环境要求、安装步骤、启动方法、测试账号',
    '● 版本控制：Git + SSH密钥认证',
    '（此处插入GitHub仓库截图）',
])

# ===== Slide 8: 总结与展望 =====
add_content_slide(prs, '总结与展望', [
    '● 已完成：3角色权限分流、每角色3+核心功能、ECharts数据可视化、GitHub公开仓库',
    '● 技术收获：实践了AI辅助编程、Flask全栈开发、前后端协作、权限系统设计',
    '● 未来方向：云服务器部署（公网访问）、匿名评价、NLP情感分析、SSO对接',
    '● 核心体会：AI助手提效显著，但系统设计能力和质量把控仍是开发者的核心素养',
])

# ===== Slide 9: 结尾 =====
add_title_slide(prs, '谢谢！', 'GitHub: https://github.com/cty23/project')

# 保存
output_path = r'C:\Users\cty27\.openclaw\workspace\course-evaluation-system\课程评价系统-演示.pptx'
prs.save(output_path)
print(f'PPT已生成：{output_path}')
